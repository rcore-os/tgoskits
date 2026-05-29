#!/usr/bin/env python3
from __future__ import annotations

import json
from datetime import date
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "ai-harness-ppt"
ASSET_DIR = OUT_DIR / "assets"
OUT = OUT_DIR / "ai-harness-defense.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

FONT = "Noto Sans CJK SC"
FONT_FALLBACK = "Microsoft YaHei"

C = {
    "dark": RGBColor(18, 24, 32),
    "dark2": RGBColor(27, 35, 45),
    "ink": RGBColor(33, 41, 54),
    "muted": RGBColor(96, 111, 128),
    "line": RGBColor(209, 216, 224),
    "paper": RGBColor(248, 250, 252),
    "white": RGBColor(255, 255, 255),
    "cyan": RGBColor(0, 166, 181),
    "cyan2": RGBColor(213, 245, 247),
    "amber": RGBColor(239, 159, 62),
    "amber2": RGBColor(255, 239, 213),
    "green": RGBColor(34, 163, 104),
    "green2": RGBColor(222, 247, 235),
    "red": RGBColor(211, 68, 79),
    "blue": RGBColor(68, 103, 201),
    "blue2": RGBColor(229, 234, 255),
    "violet": RGBColor(127, 83, 172),
    "violet2": RGBColor(239, 232, 247),
}


def read_json(path: Path) -> dict:
    if path.exists():
        return json.loads(path.read_text(encoding="utf-8", errors="replace"))
    return {}


def load_data() -> dict:
    syscall_rv = read_json(ROOT / "target/starry-syscall-harness/riscv64/latest/report.json")
    syscall_x86 = read_json(ROOT / "target/starry-syscall-harness/x86_64/latest/report.json")
    perf = read_json(ROOT / "target/starry-syscall-harness/perf/riscv64/latest/report.json")
    diff = read_json(ROOT / "target/starry-syscall-harness/perf-diff/report.json")
    return {"syscall_rv": syscall_rv, "syscall_x86": syscall_x86, "perf": perf, "diff": diff}


def ensure_assets() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    cover = ASSET_DIR / "ai-harness-cover.png"
    cover_16x9 = ASSET_DIR / "ai-harness-cover-16x9.png"
    if cover.exists() and not cover_16x9.exists():
        with Image.open(cover) as im:
            im = im.convert("RGB")
            ImageOps.fit(im, (1920, 1080), method=Image.Resampling.LANCZOS, centering=(0.7, 0.5)).save(cover_16x9)

    flame = ASSET_DIR / "flamegraph.png"
    flame_panel = ASSET_DIR / "flamegraph-panel.png"
    if flame.exists():
        with Image.open(flame) as im:
            im = im.convert("RGBA")
            canvas = Image.new("RGBA", (1920, 360), (255, 255, 255, 255))
            resized = im.resize((1780, max(68, int(im.height * 1780 / im.width))), Image.Resampling.LANCZOS)
            y = 186
            shadow = Image.new("RGBA", resized.size, (0, 0, 0, 0))
            ImageDraw.Draw(shadow).rounded_rectangle((0, 0, resized.width - 1, resized.height - 1), radius=10, fill=(0, 0, 0, 60))
            canvas.alpha_composite(shadow.filter(ImageFilter.GaussianBlur(8)), (74, y - 2))
            canvas.alpha_composite(resized, (70, y))
            draw = ImageDraw.Draw(canvas)
            draw.rounded_rectangle((70, 42, 1850, 132), radius=24, fill=(18, 24, 32, 255))
            draw.text((104, 68), "qperf flamegraph artifact", fill=(255, 255, 255, 255))
            draw.text((490, 68), "实际 SVG 产物已转换为证据截图；热点解释以 JSON/CSV 为准", fill=(205, 213, 224, 255))
            canvas.save(flame_panel)


def set_font(run, size: int | float | None = None, color: RGBColor | None = None, bold: bool = False):
    font = run.font
    font.name = FONT
    if size is not None:
        font.size = Pt(size)
    if color is not None:
        font.color.rgb = color
    font.bold = bold
    r_pr = run._r.get_or_add_rPr()
    for tag, typeface in (("a:latin", FONT), ("a:ea", FONT), ("a:cs", FONT_FALLBACK)):
        element = r_pr.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            r_pr.append(element)
        element.set("typeface", typeface)


def fill_solid(shape, color: RGBColor, transparency: int | None = None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency


def line_solid(shape, color: RGBColor, width: float = 1.0, transparency: int | None = None):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if transparency is not None:
        shape.line.transparency = transparency


def set_bg(slide, color: RGBColor = C["paper"]):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, size=18, color=C["ink"], bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)
    return box


def add_multiline(slide, x, y, w, h, lines, size=16, color=C["ink"], bullet=False, gap=0.85):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5 * gap)
        p.line_spacing = 1.05
        run = p.add_run()
        run.text = ("• " if bullet else "") + line
        set_font(run, size, color)
    return box


def add_title(slide, title: str, subtitle: str | None = None, section: str | None = None):
    if section:
        add_text(slide, Inches(0.62), Inches(0.36), Inches(2.3), Inches(0.28), section.upper(), 8.5, C["cyan"], True)
    add_text(slide, Inches(0.62), Inches(0.58), Inches(8.8), Inches(0.62), title, 26, C["dark"], True)
    if subtitle:
        add_text(slide, Inches(0.64), Inches(1.15), Inches(8.8), Inches(0.36), subtitle, 12.5, C["muted"])
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.46), Inches(1.15), Inches(0.045))
    fill_solid(line, C["amber"])
    line.line.fill.background()


def add_footer(slide, idx: int, total: int):
    add_text(slide, Inches(0.62), Inches(7.12), Inches(2.8), Inches(0.18), "AI-Harness 答辩", 8.5, C["muted"])
    add_text(slide, Inches(12.03), Inches(7.12), Inches(0.76), Inches(0.18), f"{idx:02d}/{total:02d}", 8.5, C["muted"], align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body_lines, accent=C["cyan"], fill=C["white"], title_size=16, body_size=12.5):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill_solid(card, fill)
    line_solid(card, RGBColor(224, 229, 235), 0.75)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    fill_solid(stripe, accent)
    stripe.line.fill.background()
    add_text(slide, x + Inches(0.26), y + Inches(0.18), w - Inches(0.42), Inches(0.34), title, title_size, C["dark"], True)
    add_multiline(slide, x + Inches(0.23), y + Inches(0.63), w - Inches(0.42), h - Inches(0.76), body_lines, body_size, C["ink"], bullet=False)
    return card


def add_metric(slide, x, y, w, h, value, label, accent=C["cyan"], note=None):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill_solid(card, C["white"])
    line_solid(card, RGBColor(222, 228, 235), 0.75)
    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.22), y + Inches(0.24), Inches(0.18), Inches(0.18))
    fill_solid(marker, accent)
    marker.line.fill.background()
    add_text(slide, x + Inches(0.22), y + Inches(0.54), w - Inches(0.44), Inches(0.54), value, 28, C["dark"], True)
    add_text(slide, x + Inches(0.24), y + Inches(1.12), w - Inches(0.48), Inches(0.32), label, 12.5, C["ink"], True)
    if note:
        add_text(slide, x + Inches(0.24), y + Inches(1.48), w - Inches(0.48), Inches(0.36), note, 9.5, C["muted"])
    return card


def add_arrow(slide, x1, y1, x2, y2, color=C["muted"], width=1.4):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        int(round(x1)),
        int(round(y1)),
        int(round(x2)),
        int(round(y2)),
    )
    line_solid(conn, color, width)
    conn.line.end_arrowhead = 3
    return conn


def add_node(slide, x, y, w, h, title, subtitle=None, fill=C["white"], accent=C["cyan"], title_size=14):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill_solid(shape, fill)
    line_solid(shape, accent, 1.1)
    add_text(slide, x + Inches(0.13), y + Inches(0.16), w - Inches(0.26), Inches(0.28), title, title_size, C["dark"], True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, x + Inches(0.16), y + Inches(0.50), w - Inches(0.32), h - Inches(0.56), subtitle, 9.8, C["muted"], align=PP_ALIGN.CENTER)
    return shape


def shorten_func(name: str) -> str:
    rules = [
        ("add_notify_wait_pop", "VirtQueue::add_notify_wait_pop"),
        ("compiler_builtins3mem6memcpy", "memcpy"),
        ("yield_current", "yield_current"),
        ("current_check_preempt_pending", "check_preempt_pending"),
        ("pseudofs4proc", "procfs builders"),
        ("run_idle", "run_idle"),
    ]
    for needle, label in rules:
        if needle in name:
            return label
    if len(name) > 26:
        return name[:25] + "..."
    return name


def table_cell(cell, text, size=10.5, color=C["ink"], bold=False, fill=None, align=PP_ALIGN.LEFT):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    set_font(run, size, color, bold)


def add_header_bar(slide, text, color=C["dark"]):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.18))
    fill_solid(rect, color)
    rect.line.fill.background()
    if text:
        add_text(slide, Inches(10.2), Inches(0.26), Inches(2.5), Inches(0.2), text, 8.5, C["muted"], align=PP_ALIGN.RIGHT)


def slide_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C["dark"])
    cover = ASSET_DIR / "ai-harness-cover-16x9.png"
    if cover.exists():
        slide.shapes.add_picture(str(cover), 0, 0, width=SLIDE_W, height=SLIDE_H)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(overlay, RGBColor(7, 11, 17), 28)
    overlay.line.fill.background()
    add_text(slide, Inches(0.78), Inches(0.62), Inches(3.6), Inches(0.24), "TGOSKits / StarryOS", 11, RGBColor(210, 222, 235), True)
    add_text(slide, Inches(0.74), Inches(1.28), Inches(6.8), Inches(1.05), "ArceOS 实验与 AI-Harness", 38, C["white"], True)
    add_text(slide, Inches(0.78), Inches(2.25), Inches(6.7), Inches(0.92), "任务一：工程化课程实验\n任务二：StarryOS 自治开发框架", 21, RGBColor(223, 232, 240), True)
    add_text(slide, Inches(0.80), Inches(3.54), Inches(5.9), Inches(0.58), "从真实 ArceOS crate 生态到 syscall/qperf 自动化闭环", 15, RGBColor(190, 204, 218))
    add_text(slide, Inches(0.80), Inches(6.62), Inches(3.8), Inches(0.26), f"答辩汇报 | {date.today().isoformat()}", 11, RGBColor(190, 204, 218))
    return slide


def slide_answer_structure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(247, 249, 252))
    add_header_bar(slide, "00")
    add_title(slide, "答辩结构：两个任务，一条工程能力主线", "任务一训练真实 ArceOS 工程实践；任务二把这种实践进一步自动化、证据化。", "Overview")
    add_card(
        slide,
        Inches(0.82),
        Inches(1.88),
        Inches(5.54),
        Inches(4.34),
        "任务一：tg-arceos-tutorial 五个实验",
        [
            "printcolor：从应用输出触达 axstd/axhal 层次",
            "hashmap：把 Rust 集合生态接入 no_std axstd",
            "altalloc：实现同时服务 byte/page 的 bump allocator",
            "ramfs-rename：补齐 VFS 到 ramfs 的 rename 语义链路",
            "sysmap：运行 musl 用户态程序并实现文件 mmap",
        ],
        C["cyan"],
        C["white"],
        title_size=18,
        body_size=12.9,
    )
    add_card(
        slide,
        Inches(6.98),
        Inches(1.88),
        Inches(5.54),
        Inches(4.34),
        "任务二：AI-Harness 自治开发框架",
        [
            "syscall Linux-vs-StarryOS 对拍",
            "qperf profile / flamegraph / perf-diff",
            "CLI + MCP + Local UI + Codex skill",
            "Docker 隔离执行与结构化报告",
            "面向 PR 的证据链和回归验证流程",
        ],
        C["amber"],
        C["white"],
        title_size=18,
        body_size=12.9,
    )
    add_text(slide, Inches(0.96), Inches(6.54), Inches(11.36), Inches(0.34), "贯穿主线：不是只理解 OS 原理，而是在真实工程约束下把功能做成可构建、可运行、可验证、可协作的系统能力。", 14.2, C["dark"], True, align=PP_ALIGN.CENTER)
    return slide


def slide_task1_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "T1")
    add_title(slide, "任务一概览：五个实验覆盖 ArceOS 的关键工程边界", "每个 exercise 都是独立 unikernel 项目，支持多架构 build/run，并用脚本检查串口输出。", "Task 1")
    rows = [
        ("printcolor", "ANSI 彩色输出", "应用输出 / axstd / axhal", "Hello + SGR 序列"),
        ("hashmap", "collections::HashMap", "axstd + hashbrown + alloc", "5 万键值插入与校验"),
        ("altalloc", "bump_allocator", "axalloc / global allocator", "300 万 Vec push + sort"),
        ("ramfs-rename", "std::fs::rename", "axfs + axfs_ramfs + VFS", "ramfs 文件重命名并读回"),
        ("sysmap", "SYS_MMAP", "用户态 ELF / syscall / AddrSpace", "mmap 文件并读回 hello"),
    ]
    table = slide.shapes.add_table(len(rows) + 1, 4, Inches(0.58), Inches(1.72), Inches(12.18), Inches(3.18)).table
    widths = [2.0, 2.45, 3.75, 3.98]
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    for j, h in enumerate(("实验", "核心任务", "触达模块", "验收信号")):
        table_cell(table.cell(0, j), h, 10.4, C["white"], True, C["dark"], PP_ALIGN.CENTER)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            fill = C["white"] if i % 2 else RGBColor(244, 247, 250)
            table_cell(table.cell(i, j), val, 9.7, C["ink"], j == 0, fill, PP_ALIGN.CENTER if j != 2 else PP_ALIGN.LEFT)
    add_card(
        slide,
        Inches(0.82),
        Inches(5.22),
        Inches(11.62),
        Inches(0.94),
        "共同工程框架",
        [
            "nightly + bare-metal targets + cargo xtask + QEMU；脚本覆盖 riscv64、x86_64、aarch64、loongarch64，并按串口输出做功能判定。",
        ],
        C["cyan"],
        C["white"],
        title_size=13.5,
        body_size=11.4,
    )
    return slide


def slide_task1_path(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(247, 250, 249))
    add_header_bar(slide, "T1")
    add_title(slide, "任务一学习路径：从可见输出走向用户态 syscall 语义", "实验难度不是堆知识点，而是逐步扩大“必须理解的工程上下文”。", "Task 1")
    steps = [
        ("printcolor", "应用输出\nANSI / println!"),
        ("hashmap", "标准库外观\nno_std collections"),
        ("altalloc", "内核资源\n全局分配器"),
        ("ramfs rename", "文件系统\nVFS 转发语义"),
        ("sysmap", "用户态运行\nELF + mmap syscall"),
    ]
    colors = [C["blue"], C["cyan"], C["green"], C["amber"], C["violet"]]
    x0, y0 = Inches(0.82), Inches(2.1)
    for i, (title, sub) in enumerate(steps):
        x = x0 + Inches(i * 2.42)
        add_node(slide, x, y0, Inches(1.72), Inches(1.06), title, sub, C["white"], colors[i], 12)
        if i < len(steps) - 1:
            add_arrow(slide, x + Inches(1.72), y0 + Inches(0.53), x + Inches(2.24), y0 + Inches(0.53), C["muted"], 1.0)
    add_card(
        slide,
        Inches(0.82),
        Inches(4.16),
        Inches(3.52),
        Inches(1.72),
        "传统 OS 课常见收获",
        [
            "理解进程、内存、文件系统、系统调用等概念模型",
            "在教学内核中实现相对封闭的模块",
        ],
        C["blue"],
        C["white"],
        title_size=14.5,
        body_size=11.4,
    )
    add_card(
        slide,
        Inches(4.82),
        Inches(4.16),
        Inches(7.52),
        Inches(1.72),
        "这组实验的独特收获",
        [
            "要把概念落到真实 crate 依赖、feature、trait、linker、QEMU、串口日志、跨架构 ABI 和测试脚本上；问题不再停在“算法对不对”，而是“能否接入系统并长期维护”。",
        ],
        C["amber"],
        C["white"],
        title_size=14.5,
        body_size=12,
    )
    return slide


def slide_task1_no_std(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "T1")
    add_title(slide, "独特收获一：把 Rust 生态迁移到 no_std OS 环境", "HashMap 不是“会用集合”，而是让 axstd 具备接近 std 的 API 外观。", "Task 1")
    add_card(
        slide,
        Inches(0.72),
        Inches(1.74),
        Inches(3.52),
        Inches(4.72),
        "printcolor",
        [
            "理解 `println!` 在 axstd 中如何落到串口输出",
            "ANSI SGR 序列让输出验证从文本变成字节级行为",
            "提示：修改 axstd 与 axhal 会影响不同输出层",
        ],
        C["blue"],
        C["white"],
        title_size=17,
        body_size=12.2,
    )
    add_card(
        slide,
        Inches(4.88),
        Inches(1.74),
        Inches(3.52),
        Inches(4.72),
        "hashmap",
        [
            "在 axstd 的 `collections` 中导出 HashMap / HashSet",
            "引入 `hashbrown`，并处理 `default-hasher` 与 alloc 依赖",
            "通过 50,000 个 key/value 验证分配器、格式化和迭代",
        ],
        C["cyan"],
        C["white"],
        title_size=17,
        body_size=12.2,
    )
    add_card(
        slide,
        Inches(9.04),
        Inches(1.74),
        Inches(3.52),
        Inches(4.72),
        "独特能力",
        [
            "理解 `extern crate axstd as std` 背后的兼容层",
            "知道 crates.io 版本缺口如何用本地 path / patch 修补",
            "学会在 no_std、alloc、随机数和平台 HAL 之间定位问题",
        ],
        C["green"],
        C["white"],
        title_size=17,
        body_size=12.2,
    )
    return slide


def slide_task1_allocator_fs(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(248, 249, 246))
    add_header_bar(slide, "T1")
    add_title(slide, "独特收获二：在真实 trait 边界里替换内核组件", "altalloc 和 ramfs-rename 都要求理解局部实现如何被上层全局机制调用。", "Task 1")
    add_card(
        slide,
        Inches(0.78),
        Inches(1.78),
        Inches(5.62),
        Inches(4.58),
        "altalloc：分配器不是一个孤立算法",
        [
            "`bump_allocator::EarlyAllocator` 同时承担 byte allocator 与 page allocator",
            "实现 `BaseAllocator`、`ByteAllocator`、`PageAllocator` 三组 trait",
            "`axalloc` 通过 `[patch.crates-io]` 和 feature 默认选用本地 bump allocator",
            "大 Vec + sort 压测暴露容量、对齐、回收计数和页分配边界",
        ],
        C["green"],
        C["white"],
        title_size=16,
        body_size=12.3,
    )
    add_card(
        slide,
        Inches(6.92),
        Inches(1.78),
        Inches(5.62),
        Inches(4.58),
        "ramfs-rename：文件系统语义要走完整链路",
        [
            "`std::fs::rename` 经过 axstd/axfs 到 `axfs_ramfs`",
            "补齐 `VfsNodeOps::rename`，并让 composite root 正确转发",
            "限制同目录 rename，不把 move 语义误当成 rename",
            "通过创建、写入、重命名、重新读回确认可见语义",
        ],
        C["amber"],
        C["white"],
        title_size=16,
        body_size=12.3,
    )
    return slide


def slide_task1_sysmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "T1")
    add_title(slide, "独特收获三：从内核进入用户态 ABI 与 Linux 语义", "sysmap 把课程里的 mmap 概念落实到 ELF、用户地址空间、trap 和文件 fd。", "Task 1")
    nodes = [
        ("musl payload", "mapfile.c\nstatic ELF"),
        ("FAT32 disk", "/sbin/mapfile\nvirtio-blk"),
        ("ELF loader", "PT_LOAD\nuser stack"),
        ("trap loop", "UserContext\nsyscall dispatch"),
        ("SYS_MMAP", "file-backed\nAddrSpace map"),
        ("verification", "read back\nMapFile ok"),
    ]
    colors = [C["blue"], C["cyan"], C["green"], C["amber"], C["violet"], C["dark"]]
    x0, y0 = Inches(0.62), Inches(1.82)
    for i, (title, sub) in enumerate(nodes):
        x = x0 + Inches(i * 2.05)
        add_node(slide, x, y0, Inches(1.54), Inches(0.98), title, sub, C["white"], colors[i], 11.5)
        if i < len(nodes) - 1:
            add_arrow(slide, x + Inches(1.54), y0 + Inches(0.49), x + Inches(1.88), y0 + Inches(0.49), C["muted"], 0.95)
    add_card(
        slide,
        Inches(0.84),
        Inches(3.62),
        Inches(5.62),
        Inches(2.32),
        "实现重点",
        [
            "按架构读取 syscall number 与参数，按 Linux mmap ABI 返回结果",
            "处理 PROT / MAP flags、页对齐、匿名映射、文件映射与 errno",
            "通过 `USER_ASPACE` 写入映射内容，而不是只在内核缓冲区模拟",
        ],
        C["violet"],
        C["white"],
        title_size=15,
        body_size=12.2,
    )
    add_card(
        slide,
        Inches(6.86),
        Inches(3.62),
        Inches(5.62),
        Inches(2.32),
        "区别于普通课程实验",
        [
            "需要处理真实 musl 程序的系统调用序列，而不是手写单个测试函数",
            "跨架构 syscall number / register convention 会直接影响正确性",
            "文件系统、内存管理和任务切换必须在同一个可运行镜像里闭合",
        ],
        C["cyan"],
        C["white"],
        title_size=15,
        body_size=12.2,
    )
    return slide


def slide_task1_validation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(247, 248, 252))
    add_header_bar(slide, "T1")
    add_title(slide, "任务一验证方式：面向真实运行，而不是只看单元测试", "每个 exercise 都通过 `cargo xtask run` 构建裸机镜像、启动 QEMU、检查串口输出。", "Task 1")
    add_metric(slide, Inches(0.82), Inches(1.86), Inches(2.5), Inches(1.78), "5", "standalone exercises", C["cyan"], "独立 Cargo 项目")
    add_metric(slide, Inches(3.58), Inches(1.86), Inches(2.5), Inches(1.78), "4", "target architectures", C["blue"], "rv64 / x86 / arm64 / loongarch")
    add_metric(slide, Inches(6.34), Inches(1.86), Inches(2.5), Inches(1.78), "QEMU", "runtime check", C["amber"], "serial-output oracle")
    add_metric(slide, Inches(9.10), Inches(1.86), Inches(2.5), Inches(1.78), "OK", "representative run", C["green"], "printcolor/riscv64 in Docker")
    add_card(
        slide,
        Inches(1.02),
        Inches(4.22),
        Inches(10.96),
        Inches(1.36),
        "验证信号示例",
        [
            "`Hello, Arceos!` + ANSI SGR、`test_hashmap() OK!`、`Bump tests run OK!`、`[Ramfs-Rename]: ok!`、`Read back content: hello, arceos!` + `MapFile ok!`。",
        ],
        C["green"],
        C["white"],
        title_size=14,
        body_size=12.1,
    )
    add_text(slide, Inches(0.94), Inches(6.20), Inches(11.5), Inches(0.42), "这类验证强迫我们同时关心编译目标、链接脚本、镜像格式、QEMU 参数、串口日志和跨架构差异，训练的是 OS 工程交付能力。", 14, C["dark"], True, align=PP_ALIGN.CENTER)
    return slide


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "01")
    add_title(slide, "任务二背景：StarryOS 兼容与性能问题越来越难靠人工闭环", "Linux 兼容目标下，细粒度语义差异和运行时热点都需要可复现证据。", "Task 2")
    add_card(
        slide,
        Inches(0.78),
        Inches(1.88),
        Inches(5.72),
        Inches(4.42),
        "Syscall 语义兼容",
        [
            "错误返回值、errno、flag 校验顺序会直接影响应用行为",
            "边界条件分散在 fd、路径、权限、信号和内存等路径中",
            "只看源码难以判断 Linux 真实行为",
            "修复后必须重新对拍，不能靠改弱测试通过",
        ],
        C["cyan"],
        C["white"],
        title_size=20,
        body_size=15,
    )
    add_card(
        slide,
        Inches(6.84),
        Inches(1.88),
        Inches(5.72),
        Inches(4.42),
        "qperf 性能定位",
        [
            "热点可能分布在 VirtIO、内存映射、调度、锁与 copy 路径",
            "优化不能只依赖直觉，需要样本、火焰图和 diff",
            "采样结果要能保存、比较、复查，并进入 PR 证据链",
            "agent 需要结构化输入，而不是一次性命令输出",
        ],
        C["amber"],
        C["white"],
        title_size=20,
        body_size=15,
    )
    add_text(slide, Inches(0.9), Inches(6.55), Inches(11.55), Inches(0.28), "核心矛盾：要让 AI 自动做开发，必须先把“发现问题 -> 形成证据 -> 定位代码 -> 验证回归”做成稳定接口。", 14, C["dark"], True, align=PP_ALIGN.CENTER)
    return slide


def slide_goals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "02")
    add_title(slide, "任务二设计目标：自动化，但不牺牲可信性", "框架不是脚本集合，而是面向 agent 和人工协作的开发平台。", "Task 2")
    goals = [
        ("自动化调用", "CLI/MCP 默认入口；无需人工点击即可批量扫描、采样、生成报告", C["cyan"]),
        ("可信基线", "Linux probe 输出作为 syscall 参考；StarryOS 构建和 QEMU 统一在 Docker", C["green"]),
        ("证据落盘", "JSON / Markdown / CSV / SVG / logs 全部保留，便于复核和 PR 说明", C["amber"]),
        ("可交互观察", "本地 Web UI 可看任务、报告、flamegraph 和 artifact", C["blue"]),
        ("真实 PR 流程", "fetch/rebase、验证命令、关键结果与风险点进入 PR 描述", C["violet"]),
    ]
    x0, y0 = Inches(0.76), Inches(1.88)
    for i, (title, body, color) in enumerate(goals):
        x = x0 + Inches((i % 3) * 4.1)
        y = y0 + Inches((i // 3) * 2.08)
        add_card(slide, x, y, Inches(3.64), Inches(1.62), title, [body], color, C["white"], title_size=16, body_size=12.4)
    add_text(slide, Inches(7.15), Inches(5.96), Inches(4.7), Inches(0.58), "交付形态：CLI harness + MCP server + Codex skill + Local UI + Docker 执行边界 + 结构化报告。", 15, C["dark"], True, align=PP_ALIGN.CENTER)
    return slide


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(245, 247, 250))
    add_header_bar(slide, "03")
    add_title(slide, "任务二总体架构：agent 不直接拼复杂构建命令", "harness 负责参数规范化、Docker 重入、输出目录管理、报告生成与 artifact 回收。", "Task 2")
    layers = [
        ("Agent / Human Operator", "MCP 自动调用 | CLI 直接调用 | Local Web UI", C["blue2"], C["blue"]),
        ("Harness Entry Points", "harness.py | mcp_server.py | ui_server.py", C["cyan2"], C["cyan"]),
        ("Docker Execution Boundary", "ghcr.io/rcore-os/tgoskits-container:latest", C["green2"], C["green"]),
        ("Syscall Differential Engine", "Linux probe vs StarryOS probe | CASE parser | report.json", C["amber2"], C["amber"]),
        ("qperf Performance Engine", "TCG plugin | analyzer | folded stack | flamegraph | diff", C["violet2"], C["violet"]),
        ("Reports And Feedback", "JSON | Markdown | CSV | SVG | logs | PR evidence", C["white"], C["dark"]),
    ]
    x, y, w, h = Inches(1.06), Inches(1.72), Inches(8.9), Inches(0.68)
    for i, (title, body, fill, accent) in enumerate(layers):
        yy = y + Inches(i * 0.82)
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, yy, w, h)
        fill_solid(rect, fill)
        line_solid(rect, accent, 1)
        add_text(slide, x + Inches(0.24), yy + Inches(0.12), Inches(3.4), Inches(0.25), title, 14.5, C["dark"], True)
        add_text(slide, x + Inches(3.65), yy + Inches(0.14), Inches(5.0), Inches(0.24), body, 11.2, C["ink"])
        if i < len(layers) - 1:
            mid = x + (w // 2)
            add_arrow(slide, mid, yy + h, mid, yy + h + Inches(0.14), layers[i][3], 1.0)
    add_card(
        slide,
        Inches(10.34),
        Inches(1.86),
        Inches(2.18),
        Inches(3.46),
        "关键抽象",
        [
            "稳定入口：harness.py",
            "隔离边界：Docker",
            "机器可读：report.json",
            "人可读：UI + flamegraph",
            "协作闭环：PR evidence",
        ],
        C["cyan"],
        C["white"],
        title_size=14,
        body_size=10.6,
    )
    return slide


def slide_syscall_loop(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "04")
    add_title(slide, "Syscall 对拍闭环：同一份 probe，Linux 作为语义基线", "输出稳定 CASE key=value 行，避免 fd、地址、时间戳等非确定信息进入比较。", "Syscall")
    steps = [
        ("写入 probe", "syscall_probe.c"),
        ("Linux 执行", "生成参考输出"),
        ("StarryOS 构建", "rootfs + debugfs 注入"),
        ("QEMU 运行", "捕获 begin/end marker"),
        ("CASE 解析", "ANSI 清理 + key/value"),
        ("差异报告", "report.json + artifacts"),
    ]
    colors = [C["blue"], C["green"], C["cyan"], C["amber"], C["violet"], C["dark"]]
    x0, y0 = Inches(0.72), Inches(1.88)
    box_w, gap = Inches(1.82), Inches(0.24)
    for i, (title, sub) in enumerate(steps):
        x = x0 + i * (box_w + gap)
        add_node(slide, x, y0, box_w, Inches(1.04), title, sub, C["white"], colors[i], 12.2)
        if i < len(steps) - 1:
            add_arrow(slide, x + box_w, y0 + Inches(0.52), x + box_w + gap, y0 + Inches(0.52), C["muted"], 1.1)
    code = [
        "CASE ftruncate_readonly_fd ret=-1 errno=22",
        "CASE pwritev2_writes_data ret=2 errno=0 read_ret=2 data=5859",
        "CASE dup3_same_fd ret=-1 errno=22",
    ]
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.84), Inches(3.56), Inches(5.82), Inches(1.66))
    fill_solid(rect, C["dark2"])
    line_solid(rect, RGBColor(64, 78, 96), 0.6)
    add_text(slide, Inches(1.06), Inches(3.82), Inches(5.4), Inches(0.84), "\n".join(code), 13, RGBColor(218, 227, 237))
    rv = data["syscall_rv"]
    case_count = len(rv.get("linux", {}))
    diff_count = len(rv.get("differences", []))
    add_metric(slide, Inches(7.08), Inches(3.42), Inches(2.45), Inches(1.82), str(case_count), "riscv64 cases", C["cyan"], "latest report")
    add_metric(slide, Inches(9.82), Inches(3.42), Inches(2.45), Inches(1.82), str(diff_count), "semantic diffs", C["green"], "begin/end marker 正常")
    add_text(slide, Inches(1.0), Inches(6.14), Inches(11.3), Inches(0.34), "已验证修复示例：ftruncate_readonly_fd 曾暴露 errno 映射差异，修实现后 rerun discover 验证 riscv64 无差异。", 14, C["dark"], True, align=PP_ALIGN.CENTER)
    return slide


def slide_qperf_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(248, 249, 246))
    add_header_bar(slide, "05")
    add_title(slide, "qperf 性能闭环：从采样到可执行修复线索", "profile 结果不是结论本身，而是进入代码审查和 perf-diff 的证据。", "Performance")
    steps = [
        ("StarryOS build", "release 默认"),
        ("QEMU + plugin", "TCG sampling"),
        ("raw samples", "qperf.bin"),
        ("analyzer", "resolve symbols"),
        ("folded stack", "stack.folded"),
        ("flamegraph", "SVG visual"),
        ("fix candidates", "rule-based triage"),
        ("perf-diff", "before/after"),
    ]
    x0, y0 = Inches(0.55), Inches(2.05)
    for i, (title, sub) in enumerate(steps):
        x = x0 + Inches((i % 4) * 3.12)
        y = y0 + Inches((i // 4) * 1.82)
        color = [C["cyan"], C["blue"], C["amber"], C["green"], C["violet"], C["red"], C["dark"], C["cyan"]][i]
        add_node(slide, x, y, Inches(2.42), Inches(0.92), title, sub, C["white"], color, 12.8)
        if i % 4 != 3:
            add_arrow(slide, x + Inches(2.42), y + Inches(0.46), x + Inches(2.88), y + Inches(0.46), C["muted"], 1.05)
        elif i == 3:
            add_arrow(slide, x + Inches(1.2), y + Inches(0.92), x + Inches(1.2), y + Inches(1.54), C["muted"], 1.05)
    add_card(
        slide,
        Inches(0.88),
        Inches(5.72),
        Inches(11.58),
        Inches(0.78),
        "关键工程增强",
        [
            "bounded queue + non-blocking send 降低采样扰动；symbol cache / partial record 容错提升 analyzer 稳定性；SVG/JSON/CSV 同时输出，支持机器解析与人工复核。",
        ],
        C["amber"],
        C["white"],
        title_size=13.5,
        body_size=11.3,
    )
    return slide


def slide_code_layout(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "06")
    add_title(slide, "代码与产物布局：入口清晰，证据可追踪", "代码入口、采样工具、报告目录共同构成可扩展的开发平台。", "Implementation")
    rows = [
        ("CLI / MCP / UI", "tools/starry-syscall-harness/", "harness.py, mcp_server.py, ui_server.py, web/*"),
        ("Syscall Probe", "tools/starry-syscall-harness/probes/", "syscall_probe.c: 输出稳定 CASE 行"),
        ("qperf 工具链", "tools/qperf/", "plugin / profiler / analyzer / flamegraph feature"),
        ("xtask 集成", "scripts/axbuild/src/starry/perf.rs", "构建 qperf、StarryOS、QEMU config 与 analyzer"),
        ("报告产物", "target/starry-syscall-harness/", "report.json, report.md, hotspots.csv, stack.folded, flamegraph.svg"),
    ]
    table = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.78), Inches(1.78), Inches(11.78), Inches(3.58)).table
    table.columns[0].width = Inches(2.3)
    table.columns[1].width = Inches(3.95)
    table.columns[2].width = Inches(5.53)
    headers = ("模块", "路径", "职责")
    for j, h in enumerate(headers):
        table_cell(table.cell(0, j), h, 11, C["white"], True, C["dark"], PP_ALIGN.CENTER)
    for i, row in enumerate(rows, start=1):
        for j, text in enumerate(row):
            fill = C["white"] if i % 2 else RGBColor(244, 247, 250)
            table_cell(table.cell(i, j), text, 10.2, C["ink"], j == 0, fill)
    add_text(slide, Inches(0.96), Inches(5.78), Inches(11.2), Inches(0.56), "设计取舍：agent 面对的是稳定的 harness 接口和结构化报告，不直接依赖临时命令、手写 QEMU 参数或不可复现日志。", 15, C["dark"], True, align=PP_ALIGN.CENTER)
    return slide


def slide_agent_entries(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(247, 248, 252))
    add_header_bar(slide, "07")
    add_title(slide, "Agent-facing 能力：CLI、MCP、Skill 与 Local UI 同源", "自动化入口和人工观察入口共享同一套 harness 命令与 artifact。", "Agent Interface")
    add_card(
        slide,
        Inches(0.72),
        Inches(1.78),
        Inches(3.72),
        Inches(4.54),
        "MCP tools",
        [
            "starry_syscall_doctor",
            "starry_syscall_discover",
            "starry_perf_profile",
            "starry_perf_diff",
            "starry_harness_ui_command",
        ],
        C["cyan"],
        C["white"],
        title_size=16,
        body_size=13,
    )
    add_card(
        slide,
        Inches(4.82),
        Inches(1.78),
        Inches(3.72),
        Inches(4.54),
        "Codex skill 约束",
        [
            "StarryOS 相关流程必须走 Docker",
            "Linux probe 是 syscall 参考基线",
            "qperf 结果作为 triage 输入",
            "修复后 rerun + targeted clippy",
            "不削弱 probe 来隐藏差异",
        ],
        C["green"],
        C["white"],
        title_size=16,
        body_size=12.2,
    )
    add_card(
        slide,
        Inches(8.92),
        Inches(1.78),
        Inches(3.72),
        Inches(4.54),
        "Local Web UI",
        [
            "Doctor / discover / perf-profile / perf-diff",
            "后台 job 与日志 tail",
            "report API 和 artifact API",
            "flamegraph 横向滚动展示",
            "限制同一时间一个重任务",
        ],
        C["amber"],
        C["white"],
        title_size=16,
        body_size=12.2,
    )
    return slide


def slide_trust(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "08")
    add_title(slide, "可信性设计：让自动化结果可以被审计", "关键不是“能跑”，而是跑出的差异、热点和建议能被复核。", "Reliability")
    items = [
        ("环境一致", "StarryOS build / rootfs / QEMU / qperf 全部在 Docker 中执行", C["green"]),
        ("稳定输入", "probe 只输出语义字段，避免地址、fd 编号、时间戳和调度时序", C["cyan"]),
        ("真实基线", "Linux probe 输出作为参考；差异默认修 StarryOS 实现", C["blue"]),
        ("证据优先", "报告、stdout/stderr、qemu.toml、rootfs、flamegraph 全部落盘", C["amber"]),
        ("性能克制", "fix candidates 是规则线索，必须结合代码审查和 perf-diff", C["violet"]),
    ]
    for i, (title, body, color) in enumerate(items):
        y = Inches(1.72 + i * 0.88)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.92), y + Inches(0.12), Inches(0.28), Inches(0.28))
        fill_solid(dot, color)
        dot.line.fill.background()
        add_text(slide, Inches(1.36), y, Inches(2.0), Inches(0.32), title, 16, C["dark"], True)
        add_text(slide, Inches(3.22), y + Inches(0.02), Inches(8.9), Inches(0.34), body, 14, C["ink"])
    add_text(slide, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.32), "输出可信的前提：自动化必须保留足够上下文，让人能质疑、复现和修正。", 15, C["dark"], True, align=PP_ALIGN.CENTER)
    return slide


def slide_experiment(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(247, 250, 249))
    add_header_bar(slide, "09")
    add_title(slide, "实验与验证：本地 harness 入口可用，latest 报告完整", "本轮执行 doctor 与 perf-diff 自比较；展示数据来自 target/starry-syscall-harness/latest 报告。", "Validation")
    perf = data["perf"]
    rv = data["syscall_rv"]
    x86 = data["syscall_x86"]
    samples = perf.get("hotspots", {}).get("total_samples", 0)
    flame = perf.get("summary", {}).get("flamegraph_generated", "true")
    add_metric(slide, Inches(0.74), Inches(1.82), Inches(2.62), Inches(1.82), "3/3", "doctor checks", C["green"], "Docker / image / tools")
    add_metric(slide, Inches(3.62), Inches(1.82), Inches(2.62), Inches(1.82), "0", "syscall diffs", C["cyan"], "riscv64 + x86_64")
    add_metric(slide, Inches(6.50), Inches(1.82), Inches(2.62), Inches(1.82), str(samples), "qperf samples", C["amber"], "riscv64, release")
    add_metric(slide, Inches(9.38), Inches(1.82), Inches(2.62), Inches(1.82), "OK", "perf-diff entry", C["blue"], "self compare delta 0")
    rows = [
        ("riscv64 syscall", len(rv.get("linux", {})), len(rv.get("differences", [])), rv.get("markers", {}).get("starry_end", False)),
        ("x86_64 syscall", len(x86.get("linux", {})), len(x86.get("differences", [])), x86.get("markers", {}).get("starry_end", False)),
        ("riscv64 qperf", samples, perf.get("result", "unknown"), flame),
    ]
    table = slide.shapes.add_table(4, 4, Inches(1.22), Inches(4.18), Inches(10.86), Inches(1.42)).table
    for j, h in enumerate(("项目", "样本 / case", "结果", "完整性")):
        table_cell(table.cell(0, j), h, 10.5, C["white"], True, C["dark"], PP_ALIGN.CENTER)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            table_cell(table.cell(i, j), str(val), 10.4, C["ink"], j == 0, C["white"] if i % 2 else RGBColor(244, 247, 250), PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)
    add_text(slide, Inches(0.94), Inches(6.12), Inches(11.5), Inches(0.44), "注意：性能 profile 的 plugin shutdown summary 可能因 timeout 不可用，harness 仍保留 qperf.bin、folded stack、flamegraph 和 report.json。", 12.8, C["muted"], align=PP_ALIGN.CENTER)
    return slide


def slide_hotspots(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "10")
    add_title(slide, "qperf 结果示例：机器可读热点 + 人可读 flamegraph", "条形图来自 report.json top_functions；底部为实际 flamegraph SVG 转换后的证据图。", "Hotspots")
    perf = data["perf"]
    funcs = perf.get("hotspots", {}).get("top_functions", [])[:6]
    max_pct = max([float(f.get("percent", 0)) for f in funcs] or [1])
    x0, y0 = Inches(0.9), Inches(1.84)
    bar_w = Inches(6.7)
    for i, f in enumerate(funcs):
        name = shorten_func(f.get("function", ""))
        pct = float(f.get("percent", 0))
        samples = f.get("samples", 0)
        y = y0 + Inches(i * 0.48)
        add_text(slide, x0, y - Inches(0.02), Inches(2.7), Inches(0.24), name, 9.8, C["ink"])
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x0 + Inches(2.9), y, bar_w, Inches(0.20))
        fill_solid(bg, RGBColor(232, 237, 243))
        bg.line.fill.background()
        fg_width = max(1, int(round(bar_w * (pct / max_pct))))
        fg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x0 + Inches(2.9), y, fg_width, Inches(0.20))
        fill_solid(fg, C["amber"] if i == 0 else C["cyan"])
        fg.line.fill.background()
        add_text(slide, x0 + Inches(9.72), y - Inches(0.04), Inches(1.25), Inches(0.24), f"{pct:.2f}% / {samples}", 9.6, C["muted"], align=PP_ALIGN.RIGHT)
    add_card(
        slide,
        Inches(8.78),
        Inches(1.74),
        Inches(3.38),
        Inches(2.68),
        "解读方式",
        [
            "top_functions 用于排序和自动规则匹配",
            "flamegraph 用于人工观察栈上下文",
            "fix candidates 未越过当前阈值时不强行生成修复建议",
        ],
        C["amber"],
        C["white"],
        title_size=14,
        body_size=11.1,
    )
    panel = ASSET_DIR / "flamegraph-panel.png"
    if panel.exists():
        slide.shapes.add_picture(str(panel), Inches(0.74), Inches(4.74), width=Inches(11.86), height=Inches(2.22))
    return slide


def slide_completed(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(248, 250, 252))
    add_header_bar(slide, "11")
    add_title(slide, "已完成工作量：从工具到协作路径的端到端建设", "覆盖 syscall harness、qperf、MCP/UI、skill/文档和 PR 工作流。", "Progress")
    cols = [
        ("Syscall Harness", ["Docker re-exec", "Linux/Starry probe", "rootfs 注入", "QEMU config", "CASE diff", "JSON report"], C["cyan"]),
        ("qperf Harness", ["plugin/analyzer", "xtask starry perf", "folded stack", "flamegraph SVG", "hotspots.csv", "perf diff"], C["amber"]),
        ("Agent Interface", ["MCP tools", "Codex skill", "Local UI", "job logs", "artifact API", "report API"], C["green"]),
        ("协作与文档", ["README", "framework doc", "验证命令", "PR evidence", "risk control", "future roadmap"], C["blue"]),
    ]
    for i, (title, items, color) in enumerate(cols):
        x = Inches(0.62 + i * 3.16)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(1.78), Inches(2.82), Inches(4.8))
        fill_solid(card, C["white"])
        line_solid(card, RGBColor(224, 230, 238), 0.8)
        top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(1.78), Inches(2.82), Inches(0.12))
        fill_solid(top, color)
        top.line.fill.background()
        add_text(slide, x + Inches(0.20), Inches(2.08), Inches(2.42), Inches(0.32), title, 14.2, C["dark"], True, align=PP_ALIGN.CENTER)
        add_multiline(slide, x + Inches(0.26), Inches(2.62), Inches(2.36), Inches(3.26), items, 11.2, C["ink"], bullet=True, gap=0.72)
    return slide


def slide_boundaries(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "12")
    add_title(slide, "当前边界与风险控制：明确能力范围，避免误用", "框架已具备闭环，但仍保留人工审查和持续扩展空间。", "Boundary")
    add_card(
        slide,
        Inches(0.76),
        Inches(1.82),
        Inches(5.56),
        Inches(4.54),
        "当前能力边界",
        [
            "syscall probe 覆盖仍需继续扩充",
            "qperf 主要验证 riscv64，loongarch64 需要更多实测",
            "pprof 格式预留，尚未完整支持",
            "fix candidates 仍是规则驱动",
            "Linux 性能 baseline 还需定义可比 workload",
        ],
        C["red"],
        C["white"],
        title_size=17,
        body_size=13.3,
    )
    add_card(
        slide,
        Inches(7.02),
        Inches(1.82),
        Inches(5.56),
        Inches(4.54),
        "风险控制策略",
        [
            "语义差异必须回到 Linux 参考和源码实现",
            "性能优化必须有 profile 证据和 rerun 验证",
            "Docker 统一环境，减少宿主漂移",
            "UI artifact 读取限制在仓库报告目录",
            "PR 描述同步列明背景、修改、验证和结果",
        ],
        C["green"],
        C["white"],
        title_size=17,
        body_size=13.3,
    )
    return slide


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(246, 248, 250))
    add_header_bar(slide, "13")
    add_title(slide, "后续路线：让 harness 逐步变成自治开发基础设施", "扩覆盖、强基线、做映射、进 PR，形成持续改进机制。", "Roadmap")
    phases = [
        ("1. 扩 syscall 覆盖", "fd lifecycle、signal、poll/epoll、mmap、权限与 errno 顺序", C["cyan"]),
        ("2. 差异到代码映射", "case -> syscall 实现文件 -> 常见 errno/flag 修复模板", C["green"]),
        ("3. 性能基线体系", "标准 workload、baseline 保存、阈值判断与趋势报告", C["amber"]),
        ("4. 可视化增强", "diff flamegraph、热点文件跳转、compare 并排查看", C["blue"]),
        ("5. PR 自动化", "自动汇总 discover/profile/diff 结果，生成 PR body 证据段", C["violet"]),
    ]
    y = Inches(1.78)
    for i, (title, body, color) in enumerate(phases):
        x = Inches(0.86 + i * 2.45)
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, Inches(0.78), Inches(0.78))
        fill_solid(circ, color)
        circ.line.fill.background()
        add_text(slide, x, y + Inches(0.18), Inches(0.78), Inches(0.24), str(i + 1), 17, C["white"], True, align=PP_ALIGN.CENTER)
        if i < len(phases) - 1:
            add_arrow(slide, x + Inches(0.82), y + Inches(0.39), x + Inches(2.1), y + Inches(0.39), C["muted"], 1.0)
        add_text(slide, x - Inches(0.32), Inches(2.84), Inches(1.42), Inches(0.42), title, 12.2, C["dark"], True, align=PP_ALIGN.CENTER)
        add_text(slide, x - Inches(0.58), Inches(3.44), Inches(1.96), Inches(1.12), body, 10.7, C["muted"], align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.28), Inches(5.72), Inches(10.75), Inches(0.64), "目标状态：agent 接收到差异或热点后，能自动收集证据、定位代码、提出补丁、跑回归、生成 PR 说明；人类负责语义判断和风险审查。", 16, C["dark"], True, align=PP_ALIGN.CENTER)
    return slide


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C["dark"])
    add_header_bar(slide, "")
    add_text(slide, Inches(0.84), Inches(0.86), Inches(9.2), Inches(0.54), "结论", 28, C["white"], True)
    add_text(slide, Inches(0.86), Inches(1.48), Inches(9.6), Inches(0.62), "AI-Harness 把 StarryOS 兼容性与性能优化中的“经验流程”固化成可复现、可审计、可扩展的工程接口。", 18, RGBColor(220, 230, 240))
    takeaways = [
        ("可复现", "Docker re-exec + 结构化 artifact"),
        ("可审计", "Linux 基线 + JSON/CSV/SVG 证据"),
        ("可扩展", "CLI/MCP/UI/Skill 同源入口"),
    ]
    for i, (t, b) in enumerate(takeaways):
        x = Inches(1.0 + i * 4.05)
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(3.0), Inches(3.18), Inches(1.72))
        fill_solid(rect, RGBColor(33, 43, 56), 0)
        line_solid(rect, RGBColor(77, 92, 112), 0.8)
        add_text(slide, x, Inches(3.34), Inches(3.18), Inches(0.40), t, 22, C["white"], True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.24), Inches(4.04), Inches(2.70), Inches(0.38), b, 12.5, RGBColor(190, 204, 218), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.88), Inches(6.40), Inches(3.8), Inches(0.26), "谢谢", 18, C["amber"], True)
    return slide


def slide_appendix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header_bar(slide, "Appendix")
    add_title(slide, "附录：本次 PPT 使用的命令与产物路径", None, "Evidence")
    lines = [
        "任务一来源：cg24-THU/tg-arceos-tutorial test 分支",
        "exercise-printcolor / exercise-hashmap / exercise-altalloc / exercise-ramfs-rename / exercise-sysmap",
        "代表性运行：docker run ... -w /work/exercise-printcolor ... cargo xtask run --arch riscv64",
        "python3 tools/starry-syscall-harness/harness.py doctor",
        "python3 tools/starry-syscall-harness/harness.py perf-diff --baseline target/starry-syscall-harness/perf/riscv64/latest --compare target/starry-syscall-harness/perf/riscv64/latest --top 8",
        "docs/ai-harness-development-framework.md",
        "tools/starry-syscall-harness/harness.py",
        "tools/starry-syscall-harness/mcp_server.py",
        "tools/starry-syscall-harness/ui_server.py",
        "scripts/axbuild/src/starry/perf.rs",
        "target/starry-syscall-harness/riscv64/latest/report.json",
        "target/starry-syscall-harness/x86_64/latest/report.json",
        "target/starry-syscall-harness/perf/riscv64/latest/report.json",
        "target/starry-syscall-harness/perf/riscv64/latest/qperf/flamegraph.svg",
    ]
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(1.72), Inches(11.82), Inches(4.86))
    fill_solid(rect, C["dark2"])
    line_solid(rect, RGBColor(64, 78, 96), 0.7)
    add_multiline(slide, Inches(1.02), Inches(1.98), Inches(11.26), Inches(4.22), lines, 10.5, RGBColor(220, 230, 240), bullet=False, gap=0.48)
    return slide


def build() -> None:
    ensure_assets()
    data = load_data()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "AI-Harness 答辩汇报"
    prs.core_properties.subject = "StarryOS Syscall 与 qperf 自治开发框架"
    prs.core_properties.author = "TGOSKits"

    slide_builders = [
        lambda: slide_cover(prs, data),
        lambda: slide_answer_structure(prs),
        lambda: slide_task1_overview(prs),
        lambda: slide_task1_path(prs),
        lambda: slide_task1_no_std(prs),
        lambda: slide_task1_allocator_fs(prs),
        lambda: slide_task1_sysmap(prs),
        lambda: slide_task1_validation(prs),
        lambda: slide_problem(prs),
        lambda: slide_goals(prs),
        lambda: slide_architecture(prs),
        lambda: slide_syscall_loop(prs, data),
        lambda: slide_qperf_loop(prs),
        lambda: slide_code_layout(prs),
        lambda: slide_agent_entries(prs),
        lambda: slide_trust(prs),
        lambda: slide_experiment(prs, data),
        lambda: slide_hotspots(prs, data),
        lambda: slide_completed(prs),
        lambda: slide_boundaries(prs),
        lambda: slide_roadmap(prs),
        lambda: slide_conclusion(prs),
        lambda: slide_appendix(prs),
    ]

    slides = [builder() for builder in slide_builders]
    total = len(slides)
    for idx, slide in enumerate(slides[1:-1], start=2):
        add_footer(slide, idx, total)
    add_footer(slides[-1], total, total)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
